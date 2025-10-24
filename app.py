import streamlit as st
import streamlit.components.v1 as components
from openai import OpenAI
from io import BytesIO
from typing import List
import PyPDF2
import docx
from pptx import Presentation
from typing import List, Dict, Any
import json
import pandas as pd
from bs4 import BeautifulSoup
import requests
from typing import List, Dict, Any
# ---------------------------------------------
# 1. Configurazione iniziale dell'app Streamlit
# ---------------------------------------------
st.set_page_config(page_title="FLC", layout="wide")

# === Stile personalizzato: sfondo verde Intesa e testo bianco ===
st.markdown("""
    <style>
        /* —————————————————————————— */
        /*  TEMA PRINCIPALE              */
        /* —————————————————————————— */
        .stApp {
            background-color: #007A33;
            color: white;
        }
        .block-container {
            padding: 3rem 4rem 2rem 4rem;
        }
        h1 {
            text-align: center;
            font-weight: bold;
            margin-bottom: 0.5rem;
            color: white !important;
        }
        .subtitle {
            text-align: center;
            font-size: 1.2rem;
            color: #f0f0f0;
            margin-bottom: 2rem;
        }
        .stMarkdown, .stText, label, p, div {
            color: #F0F0F0 !important;
        }
        .stButton > button {
            background-color: #004225;
            color: white;
            font-weight: bold;
            border-radius: 6px;
            transition: background 0.3s, transform 0.3s;
        }
        .stButton > button:hover {
            background-color: #006D3C;
            transform: scale(1.03);
        }

        /* —————————————————————————— */
        /*  FILE UPLOADER (CONTAINER)    */
        /* —————————————————————————— */
        section[data-testid="stFileUploader"] {
            background: transparent !important;
            padding: 0 !important;
            margin-bottom: 1.5rem !important;
        }

        /* —————————————————————————— */
        /*  FILE UPLOADER (BOX INTERNO)  */
        /* —————————————————————————— */
        section[data-testid="stFileUploader"] .css-1y4p8pa {
            max-width: 700px;                         /* larghezza massima */
            width: 100%;                              /* responsive fino a 700px */
            display: flex;
            align-items: center;
            justify-content: space-between;
            background: linear-gradient(145deg, #005c29, #008744) !important;
            border: 2px dashed #007A33 !important;
            border-radius: 12px !important;
            padding: 0.8rem 1rem !important;
            color: #ffffff !important;
            transition: background 0.2s, transform 0.2s;
        }
        section[data-testid="stFileUploader"] .css-1y4p8pa:hover {
            background: #007A33 !important;
            transform: scale(1.01);
        }
        section[data-testid="stFileUploader"] label {
            flex: 1;
            text-align: left !important;
            font-size: 0.95rem !important;
            font-weight: 600 !important;
            color: #F0F0F0 !important;
            margin: 0 !important;
        }
        section[data-testid="stFileUploader"] button {
            background-color: transparent !important;
            border: 1px solid #FFFFFF80 !important;
            border-radius: 4px !important;
            color: white !important;
        }
        section[data-testid="stFileUploader"] button:hover {
            background-color: #ffffff22 !important;
        }
    </style>
""", unsafe_allow_html=True)


st.markdown("<h1>IMP - ALPHA PROTOTYPE</h1>", unsafe_allow_html=True)
st.markdown("<p class='subtitle'>AI-powered insights from complex reports.</p>", unsafe_allow_html=True)


# ---------------------------------------------
# 2. Funzioni di utilità per estrarre testo
# ---------------------------------------------
def extract_text_from_pdf(file: BytesIO) -> str:
    reader = PyPDF2.PdfReader(file)
    text = ""
    for page in reader.pages:
        text += page.extract_text() or ""
    return text

def extract_text_from_docx(file: BytesIO) -> str:
    doc = docx.Document(file)
    return "\n".join([p.text for p in doc.paragraphs])

def extract_text_from_pptx(file: BytesIO) -> str:
    prs = Presentation(file)
    text = ""
    for slide in prs.slides:
        for shape in slide.shapes:
            if hasattr(shape, "text"):
                text += shape.text + "\n"
    return text

def extract_text(file) -> str:
    filename = file.name.lower()
    data = file.read()
    buffer = BytesIO(data)
    if filename.endswith('.pdf'):
        return extract_text_from_pdf(buffer)
    elif filename.endswith('.docx') or filename.endswith('.doc'):
        return extract_text_from_docx(buffer)
    elif filename.endswith('.pptx') or filename.endswith('.ppt'):
        return extract_text_from_pptx(buffer)
    else:
        try:
            return buffer.read().decode('utf-8')
        except:
            return ""

# ---------------------------------------------
# 3. Funzioni specifiche di analisi progetto (non usate)
# ---------------------------------------------
def analizza_obiettivi_progetto(text: str, client) -> str: ...
def analizza_stakeholder(text: str, client) -> str: ...
def analizza_milestone(text: str, client) -> str: ...
def analizza_kpi(text: str, client) -> str: ...
def analizza_rischi(text: str, client) -> str: ...

# ---------------------------------------------
# 4. Funzioni specifiche di analisi settore
# ---------------------------------------------
def analizza_tecnologia_settore(text: str, client) -> str:
    prompt = (
        "[Analisi Tecnologica] \n\n"
        "Forniscimi le seguenti informazioni presenti nei documenti: \n\n"
        " 1. Definizione delle tecnologie chiave; \n\n"
        " 2. Stato attuale delle tecnologie; \n\n"
        " 3. Potenziali vantaggi tecnologici; \n\n"
        " 4. Analisi di fattibilità e integrazione"
        f"Contenuto: {text}"
    )
    response = client.chat.completions.create(
        model="gpt-4.1-nano",
        messages=[
            {
                "role": "system",
                "content": ( 
                    
                    "Sei un agente specializzato nell'analisi tecnologica del settore spacetech.\n\n"
                    "CONSTRAINTS:\n"
                    "- Utilizza esclusivamente informazioni esplicitamente presenti nei documenti forniti. Senza allucinazioni\n\n"
                    "- Se una delle informazioni richieste non è presente nei documenti, rispondi esplicitamente con: 'Informazione non presente nei documenti forniti.'\n\n"
                    "- Non inserire mai ipotesi, inferenze o informazioni esterne al documento.\n\n"
                    "- Non generare contenuti generici o vaghi.\\n\n"
                    "- La risposta deve essere strutturata precisamente seguendo i punti enumerati nella richiesta.\n\n"
                    "- Se all'interno del documento non si fa riferimento al settore spacetech, non eseguire assolutamente le analisi della funzione analizza_mercato_settore"
                    "- Ogni informazione estratta deve essere fornita come output in lingua inglese. "
                )
            },
            {"role": "user", "content": prompt}
        ]
    )
    return response.choices[0].message.content

def analizza_mercato_settore(text: str, client) -> str:
    prompt = (
        "[Analisi di mercato] \n\n"
        "Forniscimi le seguenti informazioni presenti nei documenti: \n\n"
        " 1. Identificazione del mercato target; \n\n"
        " 2. Dimensione e opportunità di mercato; \n\n"
        " 3. Segmentazione del mercato; \n\n"
        " 4. Analisi della concorrenza; \n\n"
        f"Contenuto: {text}"
    )
    response = client.chat.completions.create(
        model="gpt-4.1-nano",
        messages=[
            {
                "role": "system",
                "content": ( 
                   
                    "Sei un agente specializzato nell'analisi di mercato del settore spacetech.\n\n"
                    "CONSTRAINTS:\n"
                    "- Utilizza esclusivamente informazioni esplicitamente presenti nei documenti forniti. Senza allucinazioni\n"
                    "- Se una delle informazioni richieste non è presente nei documenti, rispondi esplicitamente con: 'Informazione non presente nei documenti forniti.'\n\n"
                    "- Non inserire mai ipotesi, inferenze o informazioni esterne al documento.\n\n"
                    "- Non generare contenuti generici o vaghi.\n\n"
                    "- La risposta deve essere strutturata precisamente seguendo i punti enumerati nella richiesta.\n\n"
                    "- Se all'interno del documento non si fa riferimento al settore spacetech, non eseguire assolutamente le analisi della funzione analizza_mercato_settore\n\n"
                    "- Ogni informazione estratta deve essere sempre fornita come output in lingua inglese"
                )
            },
            {"role": "user", "content": prompt}
        ]
    )
    return response.choices[0].message.content


def analizza_strategia_settore(text: str, client) -> str:
    prompt = (
        "[Pianificazione strategica] \n\n"
        "Forniscimi le seguenti informazioni presenti nei documenti: \n\n"
        " 1. Definizione degli obiettivi di breve e lungo termine; \n\n"
        " 2. Sviluppo di un piano temporale; \n\n"
        " 3. Risorse e capacità di gestione dei rischi:; \n\n"
        f"Contenuto: {text}"
    )
    response = client.chat.completions.create(
        model="gpt-4.1-nano",
        messages=[
            {
                "role": "system",
                "content": ( 
                    
                    "Sei un agente specializzato nella pianificazione strategica del settore spacetech.\n\n"
                    "CONSTRAINTS:\n"
                    "- Utilizza esclusivamente informazioni esplicitamente presenti nei documenti forniti. Senza allucinazioni\n"
                    "- Se una delle informazioni richieste non è presente nei documenti, rispondi esplicitamente con: 'Informazione non presente nei documenti forniti.\n\n"
                    "- Non inserire mai ipotesi, inferenze o informazioni esterne al documento.\n\n"
                    "- Non generare contenuti generici o vaghi.\n\n"
                    "- La risposta deve essere strutturata precisamente seguendo i punti enumerati nella richiesta.\n\n"
                    "- Se all'interno del documento non si fa riferimento al settore spacetech, non eseguire assolutamente le analisi della funzione analizza_mercato_settore \n\n "
                    "Ogni informazione estratta deve essere sempre fornita come output in lingua inglese"
                )
            },
            {"role": "user", "content": prompt}
        ]
    )
    return response.choices[0].message.content

# ---------------------------------------------
# 5. Caricamento dei file
# ---------------------------------------------
st.header("1. Upload poject documents")
project_files = st.file_uploader(
    "Select one or more project files",
    type=['txt', 'pdf', 'doc', 'docx', 'ppt', 'pptx'],
    accept_multiple_files=True
)

st.header("2. Upload industry documents")
industry_files = st.file_uploader(
    "Select one or more industry files",
    type=['txt', 'pdf', 'doc', 'docx', 'ppt', 'pptx'],
    accept_multiple_files=True
)

# ---------------------------------------------
# 6. Analisi con OpenAI
# ---------------------------------------------
client = OpenAI(api_key=st.secrets["OPENAI_API_KEY"])

if st.button("Run analysis"):
    all_responses: List[str] = []

    # ⚠️ Mostra messaggio solo se sono stati caricati file progetto
    if project_files:
        all_responses.append("⚠️ Non è possibile effettuare analisi relative al progetto. L'addestramento non ha avuto luogo giacché non sono disponibili dati attinenti ai singoli progetti.")

    # ✅ Analisi documenti settore
    for file in industry_files:
        text = extract_text(file)
        all_responses.append(analizza_tecnologia_settore(text, client))
        all_responses.append(analizza_mercato_settore(text, client))
        all_responses.append(analizza_strategia_settore(text, client))

    # Output aggregato
    aggregated = "\n\n---\n\n".join(all_responses)
    st.subheader("Preliminary Results")
    st.text_area("Generated output", value=aggregated, height=300, key="initial_output")
    st.session_state["aggregated"] = aggregated

# ---------------------------------------------
# 7. Feedback e riformulazione dell'output
# ---------------------------------------------
st.header("3. Suggest edits or improvements")
feedback = st.text_area("Banking specialist feedback", height=150, key="feedback")
if st.button("Apply feedback"):
    prev = st.session_state.get("aggregated", "")
    combined_prompt = (
        f"Riformula il seguente testo in base a questo feedback:\n"
        f"Feedback: {feedback}\nTesto: {prev}"
    )
    response = client.chat.completions.create(
        model="gpt-4.1-nano",
        messages=[{"role": "user", "content": combined_prompt}]
    )
    st.subheader("Refined output")
    st.text_area("Latest output", value=response.choices[0].message.content, height=300, key="refined_output")
# ---------------------------------------------
# 8. Modulo di matching bandi Europei
# ---------------------------------------------
# Parola chiave di ricerca ("*" per tutti)
st.header("4. Search for matching calls")
keyword = st.text_input("Keyword for funding calls", "*")

import json
import requests
import pandas as pd
from bs4 import BeautifulSoup

if st.button("Find EU Calls"):

  # === Configurazione API TED ===
  api_url = "https://api.tech.ec.europa.eu/search-api/prod/rest/search"
  params = {
      "apiKey": "SEDIA",
      "text": "*",  # <-- puoi sostituire con una parola chiave specifica
      "pageSize": "50",
      "pageNumber": "1"
  }

  # === Query: solo GRANTS e aperti ===
  query = {
      "bool": {
          "must": [
              {"terms": {"type": ["1"]}},  # Solo GRANTS
              {"terms": {"status": ["31094501", "31094502"]}}  # Solo bandi aperti e in arrivo
          ]
      }
  }
  languages = ["en"]
  sort = {"field": "sortStatus", "order": "ASC"}

  # === Funzioni helper ===
  def safe_first(value, default="N/A"):
      if isinstance(value, dict):
          return next(iter(value.values())) if value else default
      elif isinstance(value, list):
          return value[0] if value else default
      elif isinstance(value, str):
          return value
      return default

  def clean_html(raw_html):
      return BeautifulSoup(raw_html, "html.parser").get_text(strip=True)

  def convert_json_url_to_web(url):
      if not url.endswith(".json"):
          raise ValueError("L'URL deve terminare con '.json'")
      topic_id = url.split("/")[-1].replace(".json", "")
      topic_id_lower = topic_id.lower()
      return (
          "https://ec.europa.eu/info/funding-tenders/opportunities/portal/screen/"
          f"opportunities/topic-details/{topic_id_lower}"
      )

  # === Scarica risultati ===
  page = 1
  results = []

  while True:
      params["pageNumber"] = str(page)
      response = requests.post(
          api_url,
          params=params,
          files={
              "query": ("blob", json.dumps(query), "application/json"),
              "languages": ("blob", json.dumps(languages), "application/json"),
              "sort": ("blob", json.dumps(sort), "application/json"),
          },
      )
      data = response.json()
      if not data.get("results"):
          break

      for result in data["results"]:
          metadata = result.get("metadata", {})
          title = safe_first(metadata.get("title"))
          raw_description = safe_first(metadata.get("descriptionByte"))
          clean_description = clean_html(raw_description)
          url = safe_first(metadata.get("url"))
          status = safe_first(metadata.get("status"))
          results.append({
              "Title": title,
              "Task Description": clean_description,
              "URL": convert_json_url_to_web(url),
              "Status": "Aperto" if status == "31094502" else "In arrivo"
          })

      if len(data["results"]) < int(params["pageSize"]):
          break
      page += 1

  df = pd.DataFrame(results)
  filtered_df = df[df["Task Description"].str.contains(keyword, case=False, na=False)] 
  urls = filtered_df["URL"].dropna().unique() 
  st.markdown("### EU calls disponibili") 
  for u in df["URL"].dropna().unique():
    st.markdown(f"- [{u}]({u})")